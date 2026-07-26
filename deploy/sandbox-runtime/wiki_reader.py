#!/opt/cocola/venv/bin/python
"""Read-only Office helpers for Cocola Wiki files materialized in a sandbox."""

from __future__ import annotations

import argparse
import csv
import json
import sys
import zipfile
from pathlib import Path

MAX_OUTPUT_BYTES = 1 << 20
MAX_OFFICE_EXPANDED_BYTES = 64 << 20
MAX_OFFICE_ARCHIVE_ENTRIES = 10_000
MAX_XLSX_RANGE_CELLS = 10_000
MAX_XLSX_ROW = 1_048_576
MAX_XLSX_COLUMN = 16_384


class _OutputLimitReached(Exception):
    pass


class _LimitedWriter:
    def __init__(self, target, max_bytes: int):
        self._target = target
        self._remaining = max_bytes

    def write(self, value: str) -> int:
        encoded = value.encode("utf-8")
        if len(encoded) <= self._remaining:
            self._remaining -= len(encoded)
            return self._target.write(value)
        prefix = encoded[: self._remaining].decode("utf-8", errors="ignore")
        if prefix:
            self._target.write(prefix)
        self._target.write(f"\n[output truncated after {MAX_OUTPUT_BYTES} bytes]\n")
        self._remaining = 0
        raise _OutputLimitReached

    def flush(self) -> None:
        self._target.flush()


def _validate_office_archive(path: Path) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_OFFICE_ARCHIVE_ENTRIES:
                raise SystemExit(
                    f"Office archive has more than {MAX_OFFICE_ARCHIVE_ENTRIES} entries"
                )
            expanded_bytes = 0
            for entry in entries:
                if entry.is_dir():
                    continue
                if (
                    entry.file_size < 0
                    or entry.file_size > MAX_OFFICE_EXPANDED_BYTES - expanded_bytes
                ):
                    raise SystemExit("Office archive expands beyond the sandbox reader limit")
                expanded_bytes += entry.file_size
    except zipfile.BadZipFile as exc:
        raise SystemExit("invalid Office archive") from exc


def _docx(path: Path) -> None:
    from docx import Document

    _validate_office_archive(path)
    document = Document(path)
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name if paragraph.style else "").lower()
        if style.startswith("heading"):
            level = "".join(char for char in style if char.isdigit()) or "1"
            print(f"{'#' * min(int(level), 6)} {text}")
        else:
            print(text)
        print()
    for table_index, table in enumerate(document.tables, 1):
        print(f"## Table {table_index}")
        rows = iter(table.rows)
        first = next(rows, None)
        if first is None:
            continue
        header = [cell.text.replace("\n", " ").strip() for cell in first.cells]
        print("| " + " | ".join(header) + " |")
        print("| " + " | ".join("---" for _ in header) + " |")
        for row in rows:
            values = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            print("| " + " | ".join(values) + " |")
        print()


def _pptx(path: Path) -> None:
    from pptx import Presentation

    _validate_office_archive(path)
    presentation = Presentation(path)
    for index, slide in enumerate(presentation.slides, 1):
        print(f"# Slide {index}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    print(text)
                    print()
            if getattr(shape, "has_table", False):
                rows = iter(shape.table.rows)
                first = next(rows, None)
                if first is not None:
                    header = [cell.text.replace("\n", " ").strip() for cell in first.cells]
                    print("| " + " | ".join(header) + " |")
                    print("| " + " | ".join("---" for _ in header) + " |")
                    for row in rows:
                        values = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                        print("| " + " | ".join(values) + " |")
                    print()
        notes_slide = getattr(slide, "notes_slide", None)
        notes_frame = getattr(notes_slide, "notes_text_frame", None)
        notes = getattr(notes_frame, "text", "").strip()
        if notes:
            print("## Speaker notes")
            print(notes)
            print()


def _xlsx_info(path: Path) -> None:
    from openpyxl import load_workbook

    _validate_office_archive(path)
    workbook = load_workbook(path, read_only=True, data_only=False)
    try:
        result = {
            "sheets": [
                {
                    "name": sheet.title,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column,
                }
                for sheet in workbook.worksheets
            ]
        }
        print(json.dumps(result, ensure_ascii=False, indent=2))
    finally:
        workbook.close()


def _xlsx_range(path: Path, sheet_name: str, cell_range: str, mode: str) -> None:
    from openpyxl import load_workbook
    from openpyxl.utils.cell import range_boundaries

    min_column, min_row, max_column, max_row = range_boundaries(cell_range)
    if (
        min_column is None
        or min_row is None
        or max_column is None
        or max_row is None
        or min_column < 1
        or min_row < 1
        or max_column > MAX_XLSX_COLUMN
        or max_row > MAX_XLSX_ROW
    ):
        raise SystemExit("range must use bounded Excel rows and columns")
    cell_count = (max_row - min_row + 1) * (max_column - min_column + 1)
    if cell_count > MAX_XLSX_RANGE_CELLS:
        raise SystemExit(f"range exceeds the {MAX_XLSX_RANGE_CELLS}-cell reader limit; split it")
    _validate_office_archive(path)
    workbook = load_workbook(path, read_only=True, data_only=mode == "cached")
    try:
        if sheet_name not in workbook.sheetnames:
            raise SystemExit(f"sheet not found: {sheet_name}")
        sheet = workbook[sheet_name]
        writer = csv.writer(sys.stdout)
        for row in sheet.iter_rows(
            min_row=min_row,
            max_row=max_row,
            min_col=min_column,
            max_col=max_column,
            values_only=True,
        ):
            writer.writerow(["" if value is None else value for value in row])
    finally:
        workbook.close()


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Read DOCX, XLSX, and PPTX Wiki files without modifying them."
    )
    commands = parser.add_subparsers(dest="command", required=True)
    for command in ("docx", "pptx", "xlsx-info"):
        command_parser = commands.add_parser(command)
        command_parser.add_argument("path", type=Path)
    range_parser = commands.add_parser("xlsx-range")
    range_parser.add_argument("path", type=Path)
    range_parser.add_argument("--sheet", required=True)
    range_parser.add_argument("--range", dest="cell_range", required=True)
    range_parser.add_argument("--mode", choices=("formula", "cached"), default="formula")
    args = parser.parse_args()
    if not args.path.is_file():
        raise SystemExit(f"file not found: {args.path}")
    original_stdout = sys.stdout
    sys.stdout = _LimitedWriter(original_stdout, MAX_OUTPUT_BYTES)
    try:
        if args.command == "docx":
            _docx(args.path)
        elif args.command == "pptx":
            _pptx(args.path)
        elif args.command == "xlsx-info":
            _xlsx_info(args.path)
        else:
            _xlsx_range(args.path, args.sheet, args.cell_range, args.mode)
    except _OutputLimitReached:
        pass
    finally:
        sys.stdout.flush()
        sys.stdout = original_stdout


if __name__ == "__main__":
    main()
